"""
PDF to PPTX Converter with Better Format Preservation
Uses pdfplumber for better text extraction with positioning
"""

import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def pdf_to_pptx_preserve_format(pdf_path, output_path=None):
    """
    Convert PDF to PPTX preserving layout and formatting as much as possible.
    
    Args:
        pdf_path (str): Path to the input PDF file
        output_path (str): Path for the output PPTX file (optional)
    
    Returns:
        str: Path to the created PPTX file
    """
    if output_path is None:
        output_path = pdf_path.rsplit('.', 1)[0] + '.pptx'
    
    print(f"Converting {pdf_path} to PowerPoint...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Open PDF
    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        
        for page_num, page in enumerate(pdf.pages):
            print(f"Processing page {page_num + 1}/{total_pages}...")
            
            # Create blank slide
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Extract text with positioning
            words = page.extract_words(
                x_tolerance=3,
                y_tolerance=3,
                keep_blank_chars=True,
                use_text_flow=True
            )
            
            if not words:
                print(f"  Warning: Page {page_num + 1} has no text")
                continue
            
            # Get page dimensions
            page_width = float(page.width)
            page_height = float(page.height)
            
            # Group words into lines based on y-position
            lines = {}
            for word in words:
                y_pos = round(float(word['top']), 1)
                if y_pos not in lines:
                    lines[y_pos] = []
                lines[y_pos].append(word)
            
            # Sort lines by y-position and words by x-position
            sorted_lines = []
            for y_pos in sorted(lines.keys()):
                line_words = sorted(lines[y_pos], key=lambda w: float(w['x0']))
                sorted_lines.append((y_pos, line_words))
            
            # Create text boxes for each line to preserve positioning
            for y_pos, line_words in sorted_lines:
                if not line_words:
                    continue
                
                # Get line boundaries
                x_start = min(float(w['x0']) for w in line_words)
                x_end = max(float(w['x1']) for w in line_words)
                
                # Convert PDF coordinates to PowerPoint inches
                left = Inches(x_start / page_width * 10)
                top = Inches(y_pos / page_height * 7.5)
                width = Inches((x_end - x_start) / page_width * 10)
                height = Inches(0.3)  # Approximate line height
                
                # Build line text
                line_text = ' '.join(w['text'] for w in line_words)
                
                # Add text box
                try:
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = line_text
                    text_frame.word_wrap = False
                    
                    # Set font size based on original size if available
                    if line_words and 'size' in line_words[0]:
                        font_size = line_words[0]['size']
                    else:
                        font_size = 12
                    
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(font_size)
                        paragraph.font.name = 'Calibri'
                        
                except Exception as e:
                    print(f"  Warning: Could not add text box: {e}")
                    continue
    
    # Save presentation
    prs.save(output_path)
    print(f"✓ Conversion complete! Saved to: {output_path}")
    print(f"  Created {total_pages} slides")
    
    return output_path

def pdf_to_pptx_simple(pdf_path, output_path=None):
    """
    Simpler conversion that maintains text flow better.
    
    Args:
        pdf_path (str): Path to the input PDF file
        output_path (str): Path for the output PPTX file (optional)
    
    Returns:
        str: Path to the created PPTX file
    """
    if output_path is None:
        output_path = pdf_path.rsplit('.', 1)[0] + '.pptx'
    
    print(f"Converting {pdf_path} to PowerPoint (simple mode)...")
    
    # Create presentation
    prs = Presentation()
    
    # Open PDF
    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        
        for page_num, page in enumerate(pdf.pages):
            print(f"Processing page {page_num + 1}/{total_pages}...")
            
            # Extract text maintaining layout
            text = page.extract_text(layout=True)
            
            if not text or not text.strip():
                print(f"  Warning: Page {page_num + 1} has no text")
                # Add blank slide
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)
                continue
            
            # Split into lines
            lines = text.split('\n')
            
            # Detect if first line looks like a title
            if lines and len(lines[0].strip()) < 100:
                title = lines[0].strip()
                content = '\n'.join(lines[1:])
                
                # Use title and content layout
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                if slide.shapes.title:
                    slide.shapes.title.text = title
                
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.text = content
                    
                    # Preserve spacing
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.name = 'Courier New'  # Monospace helps preserve alignment
            else:
                # Use blank layout with single text box
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)
                
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6.5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = 'Courier New'
    
    prs.save(output_path)
    print(f"✓ Conversion complete! Saved to: {output_path}")
    
    return output_path

def batch_convert(pdf_folder, output_folder=None, mode='simple'):
    """
    Convert all PDF files in a folder to PPTX.
    
    Args:
        pdf_folder (str): Path to folder containing PDF files
        output_folder (str): Path to output folder (optional)
        mode (str): 'simple' or 'preserve' - conversion mode
    """
    if output_folder and not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    pdf_files = [f for f in os.listdir(pdf_folder) if f.lower().endswith('.pdf')]
    
    print(f"Found {len(pdf_files)} PDF files to convert\n")
    
    convert_func = pdf_to_pptx_simple if mode == 'simple' else pdf_to_pptx_preserve_format
    
    for pdf_file in pdf_files:
        pdf_path = os.path.join(pdf_folder, pdf_file)
        
        if output_folder:
            output_path = os.path.join(output_folder, 
                                      pdf_file.rsplit('.', 1)[0] + '.pptx')
        else:
            output_path = None
        
        try:
            convert_func(pdf_path, output_path)
        except Exception as e:
            print(f"✗ Error converting {pdf_file}: {str(e)}")
        
        print()

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("PDF to PPTX Converter - Better Format Preservation")
        print("\nUsage:")
        print("  Simple mode:   python pdf_to_pptx.py input.pdf [output.pptx]")
        print("  Preserve mode: python pdf_to_pptx.py input.pdf output.pptx --preserve")
        print("  Batch:         python pdf_to_pptx.py --batch folder [output_folder] [--preserve]")
        print("\nModes:")
        print("  simple   - Better text flow, easier to edit (default)")
        print("  preserve - Maintains exact positioning, may create many text boxes")
        sys.exit(1)
    
    if sys.argv[1] == "--batch":
        if len(sys.argv) < 3:
            print("Error: Please provide input folder")
            sys.exit(1)
        
        input_folder = sys.argv[2]
        output_folder = sys.argv[3] if len(sys.argv) > 3 and sys.argv[3] != "--preserve" else None
        mode = 'preserve' if '--preserve' in sys.argv else 'simple'
        batch_convert(input_folder, output_folder, mode)
    else:
        input_pdf = sys.argv[1]
        output_pptx = None
        mode = 'simple'
        
        for arg in sys.argv[2:]:
            if arg == '--preserve':
                mode = 'preserve'
            else:
                output_pptx = arg
        
        if mode == 'preserve':
            pdf_to_pptx_preserve_format(input_pdf, output_pptx)
        else:
            pdf_to_pptx_simple(input_pdf, output_pptx)