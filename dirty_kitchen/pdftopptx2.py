from pdfminer.high_level import extract_text
from pptx import Presentation
from pptx.util import Inches

pdf_path = "In2025.pdf"
pptx_path = "Out2025.pptx"

prs = Presentation()
blank = prs.slide_layouts[6]

page = 0
while True:
    text = extract_text(pdf_path, page_numbers=[page])
    if not text or not text.strip():
        break

    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        prs.slide_width - Inches(1),
        prs.slide_height - Inches(1)
    )
    tf = tb.text_frame
    tf.text = text.strip()
    
    page += 1

prs.save(pptx_path)
print("done", pptx_path)